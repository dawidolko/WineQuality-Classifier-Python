"""
Generator prezentacji PowerPoint dla projektu Wine Quality (temat 7).
Wersja rozszerzona — wyjaśnienia przystępne dla laika.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/Users/dawio/Desktop/WineQuality-Classifier-Python")
RES = ROOT / "results"
OUT = ROOT / "Wine_Quality_Prezentacja.pptx"

NAVY = RGBColor(0x1F, 0x2A, 0x44)
WINE = RGBColor(0x7A, 0x1F, 0x3A)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
CREAM = RGBColor(0xFB, 0xF7, 0xEF)
GRAY = RGBColor(0x55, 0x5B, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEC, 0xE3, 0xD0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xB7, 0x1C, 0x1C)
BG_BOX = RGBColor(0xF4, 0xEE, 0xE0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(s, color)
    return s


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_paragraph(slide, x, y, w, h, text, size=14, color=NAVY,
                  align=PP_ALIGN.LEFT, line_spacing=1.2):
    """Akapit zwykłego tekstu — do dłuższych opisów."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=NAVY,
                bullet_color=WINE, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(3)
        r0 = p.add_run()
        r0.text = "•  "
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color
        r0.font.name = "Calibri"
        r1 = p.add_run()
        r1.text = it
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
        r1.font.name = "Calibri"
    return tb


def add_info_box(slide, x, y, w, h, title, body,
                 title_color=WINE, body_color=NAVY, bg=WHITE):
    """Pudełko z tytułem i akapitem opisu — dobre do wyjaśniania pojęć."""
    box = add_rect(slide, x, y, w, h, bg)
    box.line.color.rgb = title_color
    box.line.width = Pt(1.25)
    add_text(slide, x + Inches(0.18), y + Inches(0.1),
             w - Inches(0.3), Inches(0.4), title,
             size=14, bold=True, color=title_color)
    add_paragraph(slide, x + Inches(0.18), y + Inches(0.5),
                  w - Inches(0.3), h - Inches(0.55), body,
                  size=12, color=body_color, line_spacing=1.18)


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), NAVY)
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.06), GOLD)
    add_text(slide, Inches(0.5), Inches(0.18), SW - Inches(1), Inches(0.5),
             title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), SW - Inches(1), Inches(0.35),
                 subtitle, size=12, color=LIGHT)


def footer_bar(slide, page):
    add_rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), NAVY)
    add_text(slide, Inches(0.5), SH - Inches(0.29), Inches(8), Inches(0.3),
             "Wine Quality • Klasyfikacja wieloklasowa • Temat 7 • Dawid Olko",
             size=9, color=LIGHT)
    add_text(slide, SW - Inches(1.5), SH - Inches(0.29), Inches(1), Inches(0.3),
             f"{page}", size=9, color=GOLD, align=PP_ALIGN.RIGHT, bold=True)


def base_slide(title, subtitle=None, page=None):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, CREAM)
    header_bar(s, title, subtitle)
    if page is not None:
        footer_bar(s, page)
    return s


# ============================================================
# SLAJD 1 — TYTUŁ
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(2.4), SW, Inches(2.6), WINE)
add_rect(s, 0, Inches(2.35), SW, Inches(0.05), GOLD)
add_rect(s, 0, Inches(5.0), SW, Inches(0.05), GOLD)

add_text(s, Inches(0.8), Inches(0.5), SW - Inches(1.6), Inches(0.5),
         "PROJEKT • TEMAT 7 • TURA III • MAJ–CZERWIEC 2026",
         size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(2.7), SW - Inches(1.6), Inches(1.1),
         "Klasyfikacja jakości wina",
         size=50, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(3.7), SW - Inches(1.6), Inches(0.7),
         "Czy komputer potrafi ocenić wino tak jak ekspert?",
         size=20, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(4.35), SW - Inches(1.6), Inches(0.5),
         "Porównanie 9 modeli uczenia maszynowego z zespołem głosującym (Majority Voting)",
         size=14, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(5.4), SW - Inches(1.6), Inches(0.5),
         "Decision Tree   •   k-NN   •   Random Forest   →   VotingClassifier",
         size=17, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(6.3), SW - Inches(1.6), Inches(0.4),
         "Dawid Olko   |   zbiór: Wine Quality (Kaggle)   |   scikit-learn, pandas, plotly, streamlit",
         size=12, color=LIGHT, align=PP_ALIGN.CENTER)

# ============================================================
# SLAJD 2 — O CO CHODZI W TYM PROJEKCIE (NOWY — wprowadzenie dla laika)
# ============================================================
s = base_slide("O czym właściwie jest ten projekt?",
               "Krótkie wyjaśnienie dla osoby bez znajomości uczenia maszynowego",
               page=2)

add_paragraph(s, Inches(0.7), Inches(1.25), Inches(12), Inches(1.2),
              "Wyobraź sobie, że masz tysiąc próbek wina, a dla każdej znasz wynik 11 pomiarów "
              "chemicznych (np. ile alkoholu, jak kwaśne, ile cukru). Eksperci ocenili każde wino "
              "w skali 3–8. Pytanie brzmi: czy komputer, patrząc tylko na te liczby, potrafi "
              "zgadnąć ocenę eksperta dla nowego wina, którego jeszcze nie widział?",
              size=14)

add_text(s, Inches(0.7), Inches(2.6), Inches(12), Inches(0.4),
         "Co robimy w trzech krokach?", size=15, bold=True, color=WINE)

steps = [
    ("1", "UCZYMY",
     "Pokazujemy komputerowi tysiąc win z ocenami i pozwalamy mu znaleźć wzorce — "
     "np. „wina z większym alkoholem mają zwykle wyższe oceny”."),
    ("2", "TESTUJEMY",
     "Dajemy mu nowe wina (których nie widział wcześniej) i prosimy o ocenę. "
     "Sprawdzamy, ile razy zgadł poprawnie."),
    ("3", "PORÓWNUJEMY",
     "Sprawdzamy 9 różnych „mózgów” (algorytmów) i jeden zespół głosujący. "
     "Wybieramy najlepszy."),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(3.05 + i * 1.15)
    add_rect(s, Inches(0.7), y, Inches(0.85), Inches(0.95), WINE)
    add_text(s, Inches(0.7), y + Inches(0.18), Inches(0.85), Inches(0.6),
             num, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.6), y, Inches(11.0), Inches(0.95), WHITE)
    add_text(s, Inches(1.8), y + Inches(0.08), Inches(10.5), Inches(0.4),
             title, size=14, bold=True, color=WINE)
    add_paragraph(s, Inches(1.8), y + Inches(0.42), Inches(10.5), Inches(0.55),
                  desc, size=12, color=NAVY, line_spacing=1.15)

add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
         "Zadanie nazywa się klasyfikacją wieloklasową — wybór JEDNEJ z kilku możliwych ocen (3, 4, 5, 6, 7 lub 8).",
         size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLAJD 3 — AGENDA
# ============================================================
s = base_slide("Plan prezentacji",
               "20 slajdów — od zera do gotowych wyników",
               page=3)
agenda_left = [
    "1.  Strona tytułowa",
    "2.  O czym jest ten projekt? (dla laika)",
    "3.  Plan prezentacji",
    "4.  Temat nr 7 — co kazano nam zrobić?",
    "5.  Zbiór danych Wine Quality",
    "6.  Słownik podstawowych pojęć",
    "7.  Eksploracja: rozkład ocen wina",
    "8.  Eksploracja: korelacje cech",
    "9.  Architektura rozwiązania",
    "10. Pipeline: skalowanie + klasyfikator",
]
agenda_right = [
    "11. Trzy rodziny modeli — kto jak myśli?",
    "12. Walidacja krzyżowa 5-fold",
    "13. Zespół Majority Voting",
    "14. WYCIEK DANYCH — kluczowy slajd",
    "15. Odtwarzalność wyników",
    "16. Wykres porównujący 9 modeli",
    "17. Tabela szczegółowych wyników",
    "18. Agregaty rodzin + interpretacja",
    "19. Wygenerowane artefakty + Streamlit",
    "20. Podsumowanie i wnioski",
]
add_bullets(s, Inches(0.7), Inches(1.3), Inches(6), Inches(5.5),
            agenda_left, size=15)
add_bullets(s, Inches(6.9), Inches(1.3), Inches(6), Inches(5.5),
            agenda_right, size=15)

# ============================================================
# SLAJD 4 — TREŚĆ TEMATU
# ============================================================
s = base_slide("Treść tematu nr 7 — co dokładnie polecono nam zrobić?",
               "Cytat z polecenia + tłumaczenie na język codzienny",
               page=4)

box = add_rect(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(1.85), WHITE)
box.line.color.rgb = WINE
box.line.width = Pt(2)
add_text(s, Inches(1.0), Inches(1.35), Inches(11.3), Inches(0.4),
         "Treść polecenia (cytat oryginalny):",
         size=13, bold=True, color=WINE)
quote = ("„Bazowe modele: DecisionTree, kNN, Random Forest (każdy dla trzech wariantów "
         "hiperparametrów). Zespół predykcyjny utworzony jest za pomocą majority voting. "
         "Miary jakości: accuracy, balanced accuracy.”")
add_text(s, Inches(1.0), Inches(1.7), Inches(11.3), Inches(1.0),
         quote, size=13, color=NAVY, italic=True)
add_text(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(0.4),
         "Zbiór: Wine Quality • https://www.kaggle.com/datasets/yasserh/wine-quality-dataset",
         size=11, color=GRAY)

add_text(s, Inches(0.7), Inches(3.3), Inches(12), Inches(0.4),
         "Co to oznacza po ludzku?",
         size=15, bold=True, color=WINE)

add_paragraph(s, Inches(0.7), Inches(3.75), Inches(12), Inches(0.8),
              "Mamy zbudować trzy różne typy modeli (jak trzy różne sposoby myślenia), "
              "każdy w trzech ustawieniach (np. „proste drzewo”, „średnie drzewo”, „głębokie drzewo”). "
              "Łącznie daje to 9 modeli. Potem łączymy ich w zespół, który głosuje większością.",
              size=13)

add_text(s, Inches(0.7), Inches(4.65), Inches(12), Inches(0.4),
         "Co dokładnie projekt realizuje?",
         size=15, bold=True, color=WINE)
real = [
    "3 rodziny modeli × 3 warianty ustawień  →  9 niezależnych klasyfikatorów.",
    "Zespół: VotingClassifier z głosowaniem większościowym (każdy model oddaje 1 głos).",
    "Walidacja krzyżowa 5-fold: każdy model testowany na 5 różnych podziałach danych.",
    "Dwie miary jakości: accuracy (% trafień) i balanced_accuracy (uczciwa dla rzadkich klas).",
    "Brak przecieku danych: dane testowe nigdy nie wpływają na uczenie.",
]
add_bullets(s, Inches(0.7), Inches(5.1), Inches(12), Inches(2.2),
            real, size=13)

# ============================================================
# SLAJD 5 — ZBIÓR DANYCH
# ============================================================
s = base_slide("Zbiór danych — Wine Quality",
               "1143 próbek wina czerwonego z 11 pomiarów chemicznych",
               page=5)

add_paragraph(s, Inches(0.7), Inches(1.25), Inches(12), Inches(0.9),
              "Zbiór pochodzi z portalu Kaggle. Każdy wiersz = jedna próbka wina. "
              "Każda kolumna = jeden pomiar chemiczny (cecha) lub ocena jakości (etykieta). "
              "Naszym zadaniem jest nauczyć komputer przewidywać kolumnę „quality” "
              "na podstawie pozostałych 11 kolumn.",
              size=13)

def info_card(slide, x, y, w, h, label, value, color=WINE):
    add_rect(slide, x, y, w, h, WHITE)
    add_rect(slide, x, y, w, Inches(0.18), color)
    add_text(slide, x + Inches(0.15), y + Inches(0.28),
             w - Inches(0.3), Inches(0.5),
             value, size=26, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.15), y + Inches(0.85),
             w - Inches(0.3), Inches(0.4),
             label, size=11, color=GRAY)

info_card(s, Inches(0.7), Inches(2.3), Inches(2.85), Inches(1.3),
          "Liczba próbek wina", "1 143", WINE)
info_card(s, Inches(3.7), Inches(2.3), Inches(2.85), Inches(1.3),
          "Pomiary chemiczne", "11", WINE)
info_card(s, Inches(6.7), Inches(2.3), Inches(2.85), Inches(1.3),
          "Możliwe oceny (klasy)", "3 – 8", GOLD)
info_card(s, Inches(9.7), Inches(2.3), Inches(2.95), Inches(1.3),
          "Brakujące wartości", "0", GREEN)

add_text(s, Inches(0.7), Inches(3.85), Inches(12), Inches(0.4),
         "Cechy wejściowe (X) — co mierzymy w winie:",
         size=15, bold=True, color=WINE)
features = [
    "fixed acidity, volatile acidity, citric acid — kwasowość (kwas winowy, octowy, cytrynowy).",
    "residual sugar — cukier resztkowy (g/L); chlorides — sól; sulphates — siarczany.",
    "free / total sulfur dioxide — dwutlenek siarki (wolny i całkowity), naturalny konserwant.",
    "density — gęstość (g/cm³); pH — kwasowość pH; alcohol — zawartość alkoholu (% obj.).",
]
add_bullets(s, Inches(0.7), Inches(4.3), Inches(12), Inches(1.9),
            features, size=12)

add_text(s, Inches(0.7), Inches(6.25), Inches(12), Inches(0.4),
         "Etykieta (y): kolumna quality — średnia ocena 3 sommelierów (3 = najgorzej, 8 = najlepiej).",
         size=13, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
         "Kolumna „Id” usuwana przed nauką — to numer wiersza, nie cecha próbki.",
         size=12, color=GRAY, italic=True)

# ============================================================
# SLAJD 6 — SŁOWNIK POJĘĆ (NOWY — dla laika)
# ============================================================
s = base_slide("Słownik pojęć — co znaczą hasła używane w projekcie?",
               "Krótkie definicje przed wejściem w szczegóły",
               page=6)

terms_left = [
    ("Klasyfikacja",
     "Przyporządkowanie próbki do jednej z kilku kategorii. Tu: do jednej z ocen 3–8."),
    ("Cecha (feature)",
     "Liczba opisująca próbkę — np. zawartość alkoholu. W naszym projekcie jest ich 11."),
    ("Etykieta (label)",
     "Prawidłowa odpowiedź, której uczy się model. Tu: ocena eksperta."),
    ("Model",
     "Program komputerowy, który po nauczeniu potrafi wskazać klasę dla nowej próbki."),
    ("Hiperparametr",
     "Ustawienie modelu wybierane PRZED treningiem — np. „głębokość drzewa = 5”."),
]
terms_right = [
    ("Trening / test",
     "Trening = pokazujemy modelowi dane z odpowiedziami. Test = sprawdzamy, czy zgaduje na nowych."),
    ("Walidacja krzyżowa",
     "Powtarzanie testu 5 razy na różnych podziałach danych — wynik jest uczciwszy."),
    ("Skalowanie",
     "Zamiana liczb tak, by były w porównywalnym zakresie (np. 0–1). Pomaga niektórym modelom."),
    ("Accuracy",
     "Procent zgadniętych próbek. Prosta i intuicyjna miara skuteczności."),
    ("Wyciek danych",
     "Gdy informacja z testu „przecieka” do treningu — wyniki wtedy kłamią (są zawyżone)."),
]

col_w = Inches(6.05)
row_h = Inches(0.95)
y0 = Inches(1.3)
for i, (t, d) in enumerate(terms_left):
    y = y0 + i * (row_h + Inches(0.08))
    add_info_box(s, Inches(0.55), y, col_w, row_h, t, d)
for i, (t, d) in enumerate(terms_right):
    y = y0 + i * (row_h + Inches(0.08))
    add_info_box(s, Inches(6.75), y, col_w, row_h, t, d)

# ============================================================
# SLAJD 7 — ROZKŁAD JAKOŚCI
# ============================================================
s = base_slide("Eksploracja danych — jak wyglądają oceny wina w zbiorze?",
               "Wykres rozkładu kolumny „quality”",
               page=7)
s.shapes.add_picture(str(RES / "img_rozkald.png"),
                     Inches(0.5), Inches(1.25), width=Inches(8.3))

add_text(s, Inches(9.0), Inches(1.3), Inches(4.0), Inches(0.4),
         "Co to za wykres?",
         size=14, bold=True, color=WINE)
add_paragraph(s, Inches(9.0), Inches(1.7), Inches(4.0), Inches(1.5),
              "Słupki pokazują, ile próbek wina w naszym zbiorze ma daną ocenę. "
              "Im wyższy słupek, tym więcej takich win mamy.",
              size=11, color=NAVY)

add_text(s, Inches(9.0), Inches(3.3), Inches(4.0), Inches(0.4),
         "Najważniejsze obserwacje:",
         size=14, bold=True, color=WINE)
add_bullets(s, Inches(9.0), Inches(3.75), Inches(4.0), Inches(3),
            ["Najwięcej win ma ocenę 5 lub 6 (przeciętne).",
             "Wina bardzo dobre (7) i bardzo słabe (4) są rzadsze.",
             "Skrajne oceny 3 i 8 mają tylko po kilkanaście próbek.",
             "Zbiór jest niezbalansowany — trudno nauczyć się rzadkich klas."],
            size=11)
add_paragraph(s, Inches(9.0), Inches(6.1), Inches(4.0), Inches(1.0),
              "Dlatego oprócz „accuracy” liczymy też „balanced_accuracy” — "
              "uczciwszą miarę przy nierównych klasach.",
              size=10, color=GRAY)

# ============================================================
# SLAJD 8 — KORELACJE
# ============================================================
s = base_slide("Eksploracja — które cechy idą w parze z jakością?",
               "Macierz korelacji (współczynnik Pearsona)",
               page=8)
s.shapes.add_picture(str(RES / "img_korelacja.png"),
                     Inches(0.5), Inches(1.2), width=Inches(7.5))

add_text(s, Inches(8.2), Inches(1.3), Inches(4.7), Inches(0.4),
         "Jak czytać tę mapę kolorów?",
         size=14, bold=True, color=WINE)
add_paragraph(s, Inches(8.2), Inches(1.7), Inches(4.7), Inches(1.6),
              "Każda komórka pokazuje, jak silnie dwie cechy są ze sobą powiązane. "
              "Wartości od –1 do +1. Czerwone = razem rosną. Niebieskie = jedna rośnie, "
              "druga maleje. Białe = brak związku.",
              size=11, color=NAVY)

add_text(s, Inches(8.2), Inches(3.4), Inches(4.7), Inches(0.4),
         "Co z tego wynika dla naszego problemu?",
         size=14, bold=True, color=WINE)
add_bullets(s, Inches(8.2), Inches(3.85), Inches(4.7), Inches(3.3),
            ["alcohol — najsilniej dodatnio związany z oceną (więcej alkoholu → wyższa ocena).",
             "volatile acidity — ujemnie (octowy posmak obniża ocenę).",
             "fixed acidity i density — silnie skorelowane między sobą (gęstsze wino bywa kwaśniejsze).",
             "Wykres służy TYLKO do oglądania — nie wybieramy na jego podstawie cech (uniknięcie wycieku).",
             "Do modeli trafia komplet wszystkich 11 cech."],
            size=11)

# ============================================================
# SLAJD 9 — ARCHITEKTURA
# ============================================================
s = base_slide("Architektura rozwiązania — jak działa cały eksperyment?",
               "Pięć kroków od pliku CSV do tabel z wynikami",
               page=9)

stages = [
    ("WineQT.csv", "Wczytanie danych\nsprawdzenie braków", NAVY),
    ("Pipeline", "Skalowanie\n→ klasyfikator", WINE),
    ("5-fold CV", "Walidacja\n krzyżowa", GOLD),
    ("9 modeli\n+ Voting", "DT, kNN, RF\n+ głosowanie", WINE),
    ("Wyniki", "CSV, LaTeX,\nwykresy HTML", NAVY),
]
y = Inches(1.4)
w = Inches(2.35)
h = Inches(1.6)
for i, (title, desc, c) in enumerate(stages):
    bx = Inches(0.5 + i * 2.55)
    add_rect(s, bx, y, w, h, c)
    add_text(s, bx + Inches(0.1), y + Inches(0.18), w - Inches(0.2), Inches(0.5),
             title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = s.shapes.add_textbox(bx + Inches(0.1), y + Inches(0.75),
                              w - Inches(0.2), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(10)
        r.font.color.rgb = LIGHT
        r.font.name = "Calibri"
    if i < len(stages) - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 bx + w + Inches(0.02), y + Inches(0.6),
                                 Inches(0.13), Inches(0.4))
        fill(arr, GOLD)

add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.4),
         "Co dzieje się w każdym kroku?",
         size=15, bold=True, color=WINE)
add_paragraph(s, Inches(0.5), Inches(3.75), Inches(12.3), Inches(2.1),
              "(1) Wczytujemy plik CSV z 1143 winami i sprawdzamy, czy nie brakuje wartości. "
              "(2) Budujemy „taśmę produkcyjną” — pipeline, w którym najpierw skalujemy "
              "liczby do zakresu [0, 1], a potem podajemy je do klasyfikatora. "
              "(3) Stosujemy walidację krzyżową: dzielimy zbiór na 5 części, każdą po kolei "
              "robimy testem, a pozostałe 4 — treningiem. "
              "(4) Powtarzamy to dla 9 modeli i dla zespołu głosującego. "
              "(5) Zapisujemy wszystkie wyniki w postaci tabel (CSV, LaTeX) i wykresów (HTML, PNG).",
              size=12, line_spacing=1.25)

add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
         "Struktura kodu w repozytorium:",
         size=14, bold=True, color=WINE)
struct = [
    "src/config.py — ścieżki do plików oraz RANDOM_STATE = 42 (jeden seed na całość).",
    "src/experiment.py — definicja modeli, Pipeline, walidacja krzyżowa, eksport CSV/LaTeX.",
    "src/wykresy.py — wykresy Plotly (rozkład klas, korelacje, porównanie, agregaty).",
    "streamlit_app.py — aplikacja WWW z 4 zakładkami.",
]
add_bullets(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(1.1),
            struct, size=11)

# ============================================================
# SLAJD 10 — PIPELINE
# ============================================================
s = base_slide("Pipeline scikit-learn — fundament rozwiązania",
               "„Taśma produkcyjna”: skalowanie → klasyfikator w jednym obiekcie",
               page=10)

add_paragraph(s, Inches(0.7), Inches(1.25), Inches(12), Inches(0.9),
              "Pipeline to obiekt, który łączy w całość kilka kroków obróbki danych. "
              "Daje to dwie korzyści: kod jest czytelniejszy, a — co ważniejsze — "
              "wszystkie kroki są wykonywane w prawidłowej kolejności i tylko na "
              "danych treningowych, co eliminuje ryzyko wycieku danych.",
              size=12)

add_text(s, Inches(0.7), Inches(2.3), Inches(7), Inches(0.4),
         "Definicja w kodzie (src/experiment.py):",
         size=13, bold=True, color=WINE)

cb = add_rect(s, Inches(0.7), Inches(2.7), Inches(7.0), Inches(2.4), NAVY)
code = ("Pipeline([\n"
        '    ("preprocess", ColumnTransformer([\n'
        '        ("num", MinMaxScaler(), feature_cols)\n'
        "    ], remainder=\"drop\")),\n"
        '    ("clf", estymator),\n'
        "])")
ctx = s.shapes.add_textbox(Inches(0.85), Inches(2.85),
                           Inches(6.7), Inches(2.2)).text_frame
ctx.word_wrap = True
for i, line in enumerate(code.split("\n")):
    p = ctx.paragraphs[0] if i == 0 else ctx.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.name = "Consolas"
    r.font.size = Pt(13)
    r.font.color.rgb = LIGHT

add_text(s, Inches(8.0), Inches(2.3), Inches(4.8), Inches(0.4),
         "Co to robi krok po kroku?",
         size=13, bold=True, color=WINE)
steps = [
    "ColumnTransformer wybiera kolumny do obróbki (wszystkie 11 cech).",
    "MinMaxScaler skaluje liczby do zakresu [0, 1].",
    "Klasyfikator dostaje już ujednolicone dane.",
    "Cały Pipeline zachowuje się jak jeden model — można go uczyć i testować.",
]
add_bullets(s, Inches(8.0), Inches(2.7), Inches(4.8), Inches(2.4),
            steps, size=11)

add_text(s, Inches(0.7), Inches(5.35), Inches(12), Inches(0.4),
         "Dlaczego MinMaxScaler? — wyjaśnienie na przykładzie:",
         size=13, bold=True, color=WINE)
add_paragraph(s, Inches(0.7), Inches(5.8), Inches(12), Inches(1.5),
              "Cechy mają bardzo różne skale: alkohol jest w przedziale ok. 8–14 (% obj.), "
              "ale całkowity dwutlenek siarki sięga 6–289 mg/L. Jeśli nie ujednolicimy "
              "tych liczb, niektóre algorytmy (zwłaszcza k-NN, który mierzy odległości) "
              "potraktują dwutlenek siarki jako 30× ważniejszy niż alkohol — bo jego "
              "wartości są po prostu większe. MinMaxScaler skaluje WSZYSTKIE cechy "
              "do zakresu [0, 1], przez co stają się porównywalne.",
              size=12, line_spacing=1.25)

# ============================================================
# SLAJD 11 — MODELE BAZOWE
# ============================================================
s = base_slide("Trzy rodziny modeli — kto jak myśli?",
               "Każda rodzina to inny sposób uczenia się; każdą testujemy w 3 wariantach",
               page=11)

def family_box(slide, x, y, w, h, title, color, variants, idea, hp_explain):
    add_rect(slide, x, y, w, Inches(0.5), color)
    add_text(slide, x + Inches(0.15), y + Inches(0.07),
             w - Inches(0.3), Inches(0.4),
             title, size=15, bold=True, color=WHITE)
    body = add_rect(slide, x, y + Inches(0.5), w, h - Inches(0.5), WHITE)
    body.line.color.rgb = color
    body.line.width = Pt(1.2)
    iy = y + Inches(0.62)
    add_text(slide, x + Inches(0.15), iy, w - Inches(0.3), Inches(0.35),
             "Idea działania:", size=10, bold=True, color=color)
    add_paragraph(slide, x + Inches(0.15), iy + Inches(0.32),
                  w - Inches(0.3), Inches(1.15),
                  idea, size=10, line_spacing=1.15)
    iy2 = iy + Inches(1.55)
    add_text(slide, x + Inches(0.15), iy2, w - Inches(0.3), Inches(0.35),
             "Trzy warianty:", size=10, bold=True, color=color)
    for i, v in enumerate(variants):
        add_text(slide, x + Inches(0.25), iy2 + Inches(0.32) + i * Inches(0.32),
                 w - Inches(0.4), Inches(0.3),
                 "•  " + v, size=10, color=NAVY, font="Consolas")
    iy3 = iy2 + Inches(1.45)
    add_text(slide, x + Inches(0.15), iy3, w - Inches(0.3), Inches(0.35),
             "Co zmienia hiperparametr?", size=10, bold=True, color=color)
    add_paragraph(slide, x + Inches(0.15), iy3 + Inches(0.32),
                  w - Inches(0.3), Inches(0.9),
                  hp_explain, size=10, line_spacing=1.15)

family_box(s, Inches(0.4), Inches(1.25), Inches(4.2), Inches(6.0),
           "Decision Tree (drzewo decyzyjne)", WINE,
           ["dt_gleb_5    max_depth=5",
            "dt_gleb_15  max_depth=15",
            "dt_gleb_25  max_depth=25"],
           "Zadaje pytania typu „czy alkohol > 10%?” i schodzi gałęziami w dół, "
           "aż dotrze do liścia z odpowiedzią. Działa jak gra w 20 pytań.",
           "max_depth = ile pytań może zadać. Płytkie drzewo (5) jest proste, "
           "głębokie (25) zapamiętuje szczegóły — ryzyko przeuczenia.")
family_box(s, Inches(4.7), Inches(1.25), Inches(4.0), Inches(6.0),
           "k-NN (k najbliższych sąsiadów)", GOLD,
           ["knn_k5    n_neighbors=5",
            "knn_k11  n_neighbors=11, weights=distance",
            "knn_k21  n_neighbors=21"],
           "Dla nowego wina szuka „k” najbardziej podobnych win w bazie. "
           "Patrzy, jakie one mają oceny, i wybiera tę, która występuje najczęściej.",
           "k = ilu sąsiadów pytamy. Małe k = wrażliwe na szum, duże k = uśrednione. "
           "weights=distance = bliżsi sąsiedzi mają większy głos.")
family_box(s, Inches(8.8), Inches(1.25), Inches(4.1), Inches(6.0),
           "Random Forest (las losowy)", NAVY,
           ["rf_50    n_estimators=50",
            "rf_100  n_estimators=100",
            "rf_200  n_estimators=200"],
           "Buduje wiele drzew, każde widzi tylko fragment danych. "
           "Drzewa głosują — odpowiedź lasu to większość. Mądrość tłumu.",
           "n_estimators = ile drzew w lesie. Więcej drzew = stabilniej, "
           "ale wolniej. Każde drzewo pomaga się nawzajem korygować.")

# ============================================================
# SLAJD 12 — STRATYFIKOWANA CV
# ============================================================
s = base_slide("Walidacja krzyżowa 5-fold — uczciwy test modelu",
               "Każdy model oceniany na 5 niezależnych podziałach danych",
               page=12)

add_paragraph(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.95),
              "Po co? Gdybyśmy podzielili dane raz (np. 80% trening, 20% test), "
              "wynik mógłby być przypadkiem — zależałby od tego, które wina trafiły "
              "do testu. Walidacja krzyżowa robi to PIĘĆ razy z różnymi podziałami "
              "i podaje średnią. Wynik jest dużo bardziej wiarygodny.",
              size=12)

add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.4),
         "Wizualizacja — jak wyglądają podziały:",
         size=13, bold=True, color=WINE)

fold_y = Inches(2.75)
fold_h = Inches(0.4)
fold_w = Inches(10.5)
fold_x = Inches(1.1)
gap = Inches(0.1)
for i in range(5):
    y = fold_y + i * (fold_h + gap)
    add_text(s, Inches(0.3), y + Inches(0.07), Inches(0.7), Inches(0.4),
             f"Fold {i+1}", size=11, bold=True, color=NAVY)
    block_w = fold_w / 5
    for j in range(5):
        bx = fold_x + j * block_w
        col = WINE if j == i else NAVY
        add_rect(s, bx, y, block_w - Inches(0.05), fold_h, col)
        label = "TEST" if j == i else "TRAIN"
        add_text(s, bx, y + Inches(0.08), block_w - Inches(0.05), Inches(0.3),
                 label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.4),
         "Dlaczego stratyfikacja (StratifiedKFold)?",
         size=13, bold=True, color=WINE)
add_paragraph(s, Inches(0.7), Inches(5.85), Inches(12), Inches(1.4),
              "Pamiętaj — wina ocenione na 3 lub 8 są bardzo rzadkie. Gdybyśmy "
              "podzielili losowo, mogłoby się zdarzyć, że wszystkie wina z oceną "
              "8 trafią do treningu, a w teście nie będzie ani jednego. "
              "Stratyfikacja PILNUJE, żeby w każdej z 5 części proporcje ocen "
              "były takie same jak w całym zbiorze. Dzięki temu test zawsze widzi "
              "wszystkie klasy.",
              size=12, line_spacing=1.25)

# ============================================================
# SLAJD 13 — VOTING
# ============================================================
s = base_slide("Zespół predykcyjny — Majority Voting (głosowanie większościowe)",
               "9 modeli głosuje, wygrywa klasa z największą liczbą głosów",
               page=13)

add_paragraph(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.95),
              "Pomysł jest prosty jak głosowanie w demokracji: każdy z 9 modeli "
              "patrzy na nowe wino i mówi swoją ocenę. Liczymy głosy. Klasa "
              "z największą liczbą głosów wygrywa. To „mądrość tłumu” — "
              "zespół często myli się rzadziej niż pojedynczy członek.",
              size=12)

left_x = Inches(0.5)
top_y = Inches(2.3)

box1 = add_rect(s, left_x, top_y, Inches(4.3), Inches(4.4), WHITE)
box1.line.color.rgb = WINE
box1.line.width = Pt(1.5)
add_text(s, left_x + Inches(0.15), top_y + Inches(0.1),
         Inches(4.0), Inches(0.4),
         "Krok 1 — 9 modeli oddaje głos:",
         size=13, bold=True, color=WINE)
parts = ["dt_gleb_5  →  klasa 5",
         "dt_gleb_15 →  klasa 6",
         "dt_gleb_25 →  klasa 5",
         "knn_k5     →  klasa 6",
         "knn_k11    →  klasa 6",
         "knn_k21    →  klasa 5",
         "rf_50      →  klasa 6",
         "rf_100     →  klasa 6",
         "rf_200     →  klasa 6"]
for i, p in enumerate(parts):
    add_text(s, left_x + Inches(0.25), top_y + Inches(0.55) + i * Inches(0.4),
             Inches(3.9), Inches(0.35),
             p, size=12, color=NAVY, font="Consolas")

arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                         Inches(5.0), Inches(4.2), Inches(0.4), Inches(0.4))
fill(arr, GOLD)

center_x = Inches(5.5)
box2 = add_rect(s, center_x, top_y + Inches(0.8),
                Inches(3.7), Inches(2.5), NAVY)
add_text(s, center_x + Inches(0.15), top_y + Inches(0.95),
         Inches(3.4), Inches(0.4),
         "Krok 2 — liczenie głosów:",
         size=13, bold=True, color=GOLD)
counts = ["klasa 5:  3 głosy",
          "klasa 6:  6 głosów  ← większość!"]
for i, c in enumerate(counts):
    add_text(s, center_x + Inches(0.2), top_y + Inches(1.4) + i * Inches(0.5),
             Inches(3.3), Inches(0.4),
             c, size=14, color=WHITE, font="Consolas")
add_text(s, center_x + Inches(0.2), top_y + Inches(2.5),
         Inches(3.3), Inches(0.5),
         "Większość wygrywa.",
         size=11, italic=True, color=LIGHT)

arr2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                          Inches(9.4), Inches(4.2), Inches(0.4), Inches(0.4))
fill(arr2, GOLD)

right_x = Inches(9.9)
box3 = add_rect(s, right_x, top_y + Inches(0.8),
                Inches(3.0), Inches(2.5), WINE)
add_text(s, right_x + Inches(0.1), top_y + Inches(1.0),
         Inches(2.8), Inches(0.4),
         "Krok 3 — wynik:",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, right_x + Inches(0.1), top_y + Inches(1.7),
         Inches(2.8), Inches(0.6),
         "klasa 6",
         size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, right_x + Inches(0.1), top_y + Inches(2.55),
         Inches(2.8), Inches(0.4),
         "ostateczna ocena zespołu",
         size=10, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         "Kod:  VotingClassifier(estimators=[9 Pipeline'ów], voting=\"hard\")",
         size=12, bold=True, color=NAVY, font="Consolas", align=PP_ALIGN.CENTER)

# ============================================================
# SLAJD 14 — WYCIEK DANYCH (KLUCZOWY)
# ============================================================
s = base_slide("Wyciek danych (data leakage) — najczęstszy błąd w uczeniu maszynowym",
               "Czym jest, dlaczego psuje wyniki i jak go uniknięto w tym projekcie",
               page=14)

add_paragraph(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.85),
              "Wyciek danych występuje wtedy, gdy informacja ze zbioru testowego "
              "„przecieka” do treningu. Model wówczas sztucznie zna odpowiedzi "
              "i podaje wyniki LEPSZE niż w rzeczywistości — w produkcji okaże się "
              "dużo gorszy. Oto przykład typowego błędu:",
              size=12)

add_rect(s, Inches(0.5), Inches(2.15), Inches(6.1), Inches(2.4), WHITE)
bad_bar = add_rect(s, Inches(0.5), Inches(2.15), Inches(6.1), Inches(0.4), RED)
add_text(s, Inches(0.65), Inches(2.21), Inches(5.9), Inches(0.4),
         "ŹLE — klasyczny wyciek danych",
         size=13, bold=True, color=WHITE)
bad_code = ("scaler = MinMaxScaler()\n"
            "X = scaler.fit_transform(X)   # cały zbiór!\n"
            "scores = cross_val_score(clf, X, y, cv=5)")
ctx = s.shapes.add_textbox(Inches(0.65), Inches(2.65),
                           Inches(5.8), Inches(1.1)).text_frame
ctx.word_wrap = True
for i, line in enumerate(bad_code.split("\n")):
    p = ctx.paragraphs[0] if i == 0 else ctx.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.name = "Consolas"
    r.font.size = Pt(11)
    r.font.color.rgb = NAVY
add_paragraph(s, Inches(0.65), Inches(3.8), Inches(5.8), Inches(0.7),
              "Skaler poznaje min/max ze WSZYSTKICH danych — czyli także z testu. "
              "Test pomógł w skalowaniu, więc wynik jest zafałszowany.",
              size=10, color=RED, line_spacing=1.15)

add_rect(s, Inches(6.7), Inches(2.15), Inches(6.1), Inches(2.4), WHITE)
ok_bar = add_rect(s, Inches(6.7), Inches(2.15), Inches(6.1), Inches(0.4), GREEN)
add_text(s, Inches(6.85), Inches(2.21), Inches(5.9), Inches(0.4),
         "DOBRZE — sposób użyty w tym projekcie",
         size=13, bold=True, color=WHITE)
ok_code = ("pipe = Pipeline([\n"
           '    ("preprocess", MinMaxScaler()),\n'
           '    ("clf", estymator),\n'
           "])\n"
           "cross_validate(pipe, X, y, cv=cv, ...)")
ctx2 = s.shapes.add_textbox(Inches(6.85), Inches(2.65),
                            Inches(5.8), Inches(1.5)).text_frame
ctx2.word_wrap = True
for i, line in enumerate(ok_code.split("\n")):
    p = ctx2.paragraphs[0] if i == 0 else ctx2.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.name = "Consolas"
    r.font.size = Pt(11)
    r.font.color.rgb = NAVY
add_paragraph(s, Inches(6.85), Inches(4.05), Inches(5.8), Inches(0.5),
              "Skaler fitowany TYLKO na foldzie treningowej w każdej iteracji. Test pozostaje nietknięty.",
              size=10, color=GREEN, line_spacing=1.15)

add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
         "Pełna lista zabezpieczeń w projekcie (6 punktów):",
         size=14, bold=True, color=WINE)
checks = [
    "MinMaxScaler wewnątrz Pipeline → fit tylko na X_train w każdej z 5 foldów.",
    "Wszystkie 9 składników VotingClassifier to osobne Pipeline'y → ensemble też bez wycieku.",
    "Brak imputacji średnią — sprawdzamy isna() i blokujemy eksperyment, jeśli są braki.",
    "Macierz korelacji w EDA wyłącznie do oglądania — nie wybieramy na jej podstawie cech.",
    "Kolumna Id usunięta przed treningiem — to identyfikator wiersza, nie cecha próbki.",
    "Stratyfikacja zachowuje rozkład klas — train i test zawsze widzą wszystkie etykiety.",
]
add_bullets(s, Inches(0.7), Inches(5.15), Inches(12), Inches(2.2),
            checks, size=12, bullet_color=GREEN)

# ============================================================
# SLAJD 15 — ODTWARZALNOŚĆ
# ============================================================
s = base_slide("Odtwarzalność wyników (reproducibility)",
               "Każdy uruchamiający projekt powinien dostać IDENTYCZNE liczby",
               page=15)

add_paragraph(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.95),
              "Wiele algorytmów uczenia maszynowego używa losowości (np. losowy "
              "podział na foldy, losowe drzewa w lesie). Bez kontroli ta losowość "
              "powodowałaby, że za każdym uruchomieniem dostalibyśmy nieco inne "
              "wyniki. Tego się unika trzema technikami:",
              size=12)

cards = [
    ("Stały seed losowości", "RANDOM_STATE = 42",
     "Jedna stała w pliku src/config.py używana w każdym losowym komponencie: "
     "drzewa, Random Forest, dzielenie na foldy. Tasowanie zawsze daje ten sam wynik."),
    ("Wspólny obiekt CV", "cv = StratifiedKFold(...)",
     "Tworzymy jeden obiekt walidacji i przekazujemy go do KAŻDEGO testu. "
     "Dzięki temu wszystkie 9 modeli i zespół są oceniane na DOKŁADNIE tych "
     "samych podziałach — porównanie jest uczciwe."),
    ("Zapis wersji bibliotek", "wersje_bibliotek.txt",
     "Plik tworzony przy każdym uruchomieniu. Notuje wersje scikit-learn 1.8.0, "
     "numpy 2.4.4, pandas 2.3.3 — można odtworzyć identyczne środowisko."),
]
for i, (t, code, desc) in enumerate(cards):
    x = Inches(0.5 + i * 4.25)
    add_rect(s, x, Inches(2.3), Inches(4.0), Inches(3.3), WHITE)
    add_rect(s, x, Inches(2.3), Inches(4.0), Inches(0.5), NAVY)
    add_text(s, x + Inches(0.15), Inches(2.36), Inches(3.7), Inches(0.4),
             t, size=13, bold=True, color=WHITE)
    add_text(s, x + Inches(0.15), Inches(2.95), Inches(3.7), Inches(0.5),
             code, size=12, bold=True, color=WINE, font="Consolas")
    add_paragraph(s, x + Inches(0.15), Inches(3.5), Inches(3.7), Inches(2.0),
                  desc, size=11, color=NAVY, line_spacing=1.18)

add_text(s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.4),
         "Dlaczego shuffle=True wymaga seeda?",
         size=14, bold=True, color=WINE)
add_paragraph(s, Inches(0.7), Inches(6.3), Inches(12), Inches(1.0),
              "StratifiedKFold(shuffle=True) tasuje próbki przed podziałem na foldy. "
              "Bez random_state każde uruchomienie dawałoby inne tasowanie i inne podziały. "
              "Z seedem: identyczna kolejność tasowania → ten sam podział → "
              "ten sam wynik za każdym razem.",
              size=12, line_spacing=1.25)

# ============================================================
# SLAJD 16 — PORÓWNANIE MODELI (wykres)
# ============================================================
s = base_slide("Wyniki — porównanie 9 modeli + zespołu",
               "Średnia z 5 foldów ± odchylenie standardowe (paski błędu)",
               page=16)
s.shapes.add_picture(str(RES / "img_modele.png"),
                     Inches(0.4), Inches(1.2), width=Inches(12.5))

add_paragraph(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.9),
              "Lewy panel pokazuje accuracy (procent trafień), prawy — balanced_accuracy "
              "(uczciwa miara dla rzadkich klas). Każdy słupek = jeden model. "
              "Czarne pionowe kreski to słupki błędu — pokazują, jak bardzo wynik "
              "zmienia się między 5 foldami (im krótsze, tym stabilniejszy model).",
              size=11, color=GRAY, line_spacing=1.2)

# ============================================================
# SLAJD 17 — TABELA SZCZEGÓŁOWA
# ============================================================
s = base_slide("Wyniki szczegółowe — tabela liczb",
               "Wszystkie 10 wierszy (9 modeli bazowych + zespół Majority Voting)",
               page=17)

add_paragraph(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
              "Tabela pochodzi z pliku wyniki_szczegolowe.tex. Kolumny: średnia "
              "(_mean) i odchylenie standardowe (_std) z 5 foldów dla obu metryk.",
              size=11, color=GRAY)

table_data = [
    ("Model", "Rodzina", "Acc mean", "Acc std", "BalAcc mean", "BalAcc std"),
    ("dt_gleb_5",  "DecisionTree", "0.5660", "0.0069", "0.2738", "0.0115"),
    ("dt_gleb_15", "DecisionTree", "0.5940", "0.0161", "0.3053", "0.0205"),
    ("dt_gleb_25", "DecisionTree", "0.5625", "0.0183", "0.2717", "0.0129"),
    ("knn_k5",  "kNN", "0.5521", "0.0241", "0.2621", "0.0214"),
    ("knn_k11", "kNN", "0.6387", "0.0224", "0.3421", "0.0287"),
    ("knn_k21", "kNN", "0.5739", "0.0218", "0.2625", "0.0129"),
    ("rf_50",  "RandomForest", "0.6483", "0.0227", "0.3216", "0.0335"),
    ("rf_100", "RandomForest", "0.6571", "0.0263", "0.3446", "0.0336"),
    ("rf_200", "RandomForest", "0.6719", "0.0326", "0.3530", "0.0375"),
    ("majority_voting_9", "Ensemble", "0.6597", "0.0172", "0.3283", "0.0286"),
]
rows = len(table_data)
cols = 6
tbl = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.75),
                         Inches(12.3), Inches(4.5)).table
widths = [Inches(2.4), Inches(2.0), Inches(1.7), Inches(1.7),
          Inches(2.2), Inches(2.3)]
for i, w in enumerate(widths):
    tbl.columns[i].width = w

best_idx = 9
ens_idx = 10
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = str(table_data[r][c])
        run.font.name = "Calibri"
        if r == 0:
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            run.font.size = Pt(11)
            if r == best_idx:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = WINE
            elif r == ens_idx:
                run.font.bold = True
                run.font.color.rgb = NAVY
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
            else:
                run.font.color.rgb = NAVY
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else CREAM
        para.alignment = PP_ALIGN.LEFT if c < 2 else PP_ALIGN.CENTER

add_paragraph(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7),
              "Bordowy wiersz = najlepszy model pojedynczy (rf_200 z accuracy 0.6719). "
              "Beżowy = zespół Majority Voting. Niższe std (0.0172) oznacza, że zespół "
              "jest najBARDZIEJ stabilny między foldami.",
              size=11, color=GRAY, line_spacing=1.25)

# ============================================================
# SLAJD 18 — AGREGATY + INTERPRETACJA
# ============================================================
s = base_slide("Ranking rodzin modeli + interpretacja wyników",
               "Średnie i maksima dla każdej z 3 rodzin",
               page=18)
s.shapes.add_picture(str(RES / "img_agregaty.png"),
                     Inches(0.4), Inches(1.2), width=Inches(7.0))

add_text(s, Inches(7.7), Inches(1.25), Inches(5.4), Inches(0.4),
         "Ranking (accuracy):", size=13, bold=True, color=WINE)

ranking = [
    ("1.", "Random Forest", "0.6591", "0.6719", WINE),
    ("2.", "kNN", "0.5882", "0.6387", GOLD),
    ("3.", "Decision Tree", "0.5742", "0.5940", NAVY),
]
for i, (rank, name, mean, mx, color) in enumerate(ranking):
    y = Inches(1.7 + i * 0.95)
    add_rect(s, Inches(7.7), y, Inches(5.4), Inches(0.85), WHITE)
    add_rect(s, Inches(7.7), y, Inches(0.12), Inches(0.85), color)
    add_text(s, Inches(7.9), y + Inches(0.1), Inches(0.4), Inches(0.5),
             rank, size=20, bold=True, color=color)
    add_text(s, Inches(8.3), y + Inches(0.12), Inches(2.7), Inches(0.4),
             name, size=14, bold=True, color=NAVY)
    add_text(s, Inches(8.3), y + Inches(0.45), Inches(4.8), Inches(0.35),
             f"średnia: {mean}   |   maksimum: {mx}",
             size=11, color=GRAY)

add_text(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.4),
         "Interpretacja — co mówią nam liczby?",
         size=14, bold=True, color=WINE)
obs = [
    "Najlepszy pojedynczy model: rf_200 — accuracy 0.6719, balanced_accuracy 0.3530.",
    "Random Forest dominuje — wszystkie 3 warianty są w czołówce (0.6483, 0.6571, 0.6719).",
    "k-NN: ogromny skok przy k=11 z weights=\"distance\" (0.6387 vs 0.5521 dla k=5).",
    "Decision Tree najsłabszy — pojedyncze drzewo nie radzi sobie z niezbalansowanym zbiorem.",
    "Zespół: 0.6597 — lepszy od średniej, ale słabszy niż rf_200; za to NAJSTABILNIEJSZY (std=0.0172).",
]
add_bullets(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(2.2),
            obs, size=11)

# ============================================================
# SLAJD 19 — ARTEFAKTY + STREAMLIT
# ============================================================
s = base_slide("Wygenerowane artefakty + aplikacja Streamlit",
               "Wszystko co tworzy projekt — i jak to oglądać",
               page=19)

add_text(s, Inches(0.5), Inches(1.2), Inches(6.2), Inches(0.4),
         "Pliki w katalogu results/:",
         size=14, bold=True, color=WINE)

groups = [
    ("Tabele liczbowe (CSV)", WINE,
     "wyniki_szczegolowe.csv (10 wierszy), wyniki_agregaty_rodzin.csv (3 wiersze)."),
    ("Pliki LaTeX (do pracy magisterskiej)", GOLD,
     "wyniki_szczegolowe.tex, wyniki_agregaty_rodzin.tex — gotowe tabele do raportu."),
    ("Wykresy interaktywne (Plotly HTML)", NAVY,
     "01_rozkald_jakosci.html, 02_macierz_korelacji.html, "
     "03_porownanie_modeli.html, 04_agregaty_rodzin.html."),
    ("Metadane", GRAY,
     "wersje_bibliotek.txt — scikit-learn 1.8.0, numpy 2.4.4, pandas 2.3.3."),
]
y_cur = Inches(1.65)
for title, color, items in groups:
    add_rect(s, Inches(0.5), y_cur, Inches(6.2), Inches(0.36), color)
    add_text(s, Inches(0.65), y_cur + Inches(0.02),
             Inches(6.0), Inches(0.35),
             title, size=12, bold=True, color=WHITE)
    y_cur += Inches(0.4)
    add_paragraph(s, Inches(0.65), y_cur, Inches(6.0), Inches(0.7),
                  items, size=10, color=NAVY, line_spacing=1.2)
    y_cur += Inches(0.7)

add_text(s, Inches(7.0), Inches(1.2), Inches(6.0), Inches(0.4),
         "Aplikacja Streamlit — interaktywny podgląd wyników:",
         size=14, bold=True, color=WINE)
add_paragraph(s, Inches(7.0), Inches(1.6), Inches(6.0), Inches(0.6),
              "Uruchomienie: ./start.sh (Linux/macOS) lub start.bat (Windows). "
              "Otwiera stronę w przeglądarce z 4 zakładkami:",
              size=11, line_spacing=1.2)

tabs = [
    ("Opis projektu",
     "Cel, słownik pojęć, zgodność z wymaganiami, podgląd plików .tex."),
    ("Wymogi i metodologia",
     "Tabela: wymaganie → konkretne miejsce w kodzie."),
    ("Zbiór danych (EDA)",
     "Statystyki, wykres rozkładu, macierz korelacji, podgląd 15 wierszy."),
    ("Wyniki klasyfikacji",
     "Wykresy z CV, tabele, podgląd .tex, przyciski pobierania."),
]
y_cur = Inches(2.55)
for i, (t, d) in enumerate(tabs):
    add_rect(s, Inches(7.0), y_cur, Inches(6.0), Inches(0.95), WHITE)
    accent = WINE if i % 2 == 0 else NAVY
    add_rect(s, Inches(7.0), y_cur, Inches(0.15), Inches(0.95), accent)
    add_text(s, Inches(7.25), y_cur + Inches(0.08),
             Inches(5.7), Inches(0.4),
             f"Zakładka {i+1}: {t}",
             size=13, bold=True, color=accent)
    add_paragraph(s, Inches(7.25), y_cur + Inches(0.4),
                  Inches(5.7), Inches(0.5),
                  d, size=10, color=GRAY, line_spacing=1.15)
    y_cur += Inches(1.05)

# ============================================================
# SLAJD 20 — PODSUMOWANIE
# ============================================================
s = base_slide("Podsumowanie projektu",
               "Co udało się zrobić — najważniejsze wnioski",
               page=20)

chk_left = [
    "Zbiór z Kaggle (Wine Quality, nie Iris).",
    "Sprawdzenie braków w danych przed treningiem.",
    "MinMaxScaler wewnątrz Pipeline — bez wycieku.",
    "Stratyfikowana 5-fold cross-validation.",
    "9 modeli bazowych: 3 rodziny × 3 warianty.",
]
chk_right = [
    "Zespół Majority Voting (VotingClassifier).",
    "Metryki: accuracy + balanced_accuracy.",
    "Pandas: groupby, agg, to_latex (.tex).",
    "Odtwarzalność: RANDOM_STATE, ten sam cv.",
    "Aplikacja Streamlit z 4 zakładkami.",
]

def check_item(slide, x, y, w, text):
    add_rect(slide, x, y, Inches(0.32), Inches(0.32), GREEN)
    add_text(slide, x, y - Inches(0.02), Inches(0.32), Inches(0.4),
             "✓", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.45), y + Inches(0.04),
             w - Inches(0.5), Inches(0.4),
             text, size=12, color=NAVY)

add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.4),
         "Wszystkie wymagania polecenia — spełnione w 100%:",
         size=14, bold=True, color=WINE)

for i, t in enumerate(chk_left):
    check_item(s, Inches(0.7), Inches(1.7 + i * 0.5), Inches(6), t)
for i, t in enumerate(chk_right):
    check_item(s, Inches(6.9), Inches(1.7 + i * 0.5), Inches(6), t)

add_rect(s, Inches(0.7), Inches(4.4), Inches(12), Inches(1.6), NAVY)
add_text(s, Inches(0.85), Inches(4.5), Inches(11.7), Inches(0.4),
         "Najważniejsze liczby:",
         size=13, bold=True, color=GOLD)
add_text(s, Inches(0.85), Inches(4.95), Inches(5.8), Inches(0.5),
         "Najlepszy model:  rf_200",
         size=14, bold=True, color=WHITE)
add_text(s, Inches(0.85), Inches(5.3), Inches(5.8), Inches(0.5),
         "accuracy = 0.6719   |   balanced_accuracy = 0.3530",
         size=12, color=LIGHT, font="Consolas")
add_text(s, Inches(7.0), Inches(4.95), Inches(5.6), Inches(0.5),
         "Zespół Majority Voting:",
         size=14, bold=True, color=WHITE)
add_text(s, Inches(7.0), Inches(5.3), Inches(5.6), Inches(0.5),
         "accuracy = 0.6597   |   std = 0.0172  (najstabilniejszy)",
         size=12, color=LIGHT, font="Consolas")

add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.4),
         "Wniosek końcowy:",
         size=14, bold=True, color=WINE)
add_paragraph(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.6),
              "Na tym zbiorze danych pojedynczy Random Forest pokonał zespół głosujący — "
              "ale zespół był NAJSTABILNIEJSZY. W praktyce to ważna cecha: powtarzalność wyników.",
              size=12, line_spacing=1.2)

# ============================================================
prs.save(str(OUT))
print(f"Zapisano: {OUT}")
